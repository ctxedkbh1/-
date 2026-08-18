import re

from pptx import Presentation
from pptx.util import Pt


def _bullets(text, max_bullets=6, max_len=72):
    parts = re.split(r"[。！？]", text)
    out = []
    for p in parts:
        p = p.strip()
        if not p:
            continue
        if len(p) > max_len:
            p = p[:max_len] + "……"
        out.append(p)
        if len(out) >= max_bullets:
            break
    return out


def _sections(doc):
    sections = []
    cur = None
    for b in doc.blocks:
        if b.get("kind") == "heading1":
            cur = {"heading": b.get("text", ""), "paragraphs": []}
            sections.append(cur)
        elif b.get("kind") == "paragraph" and cur is not None:
            cur["paragraphs"].append(b.get("text", ""))
    return sections


class PptxExporter:
    def export(self, doc, path):
        prs = Presentation()

        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = doc.meta.get("title") or "论文"
        subtitle = "　".join(x for x in (
            doc.meta.get("name") and f"{doc.meta['name']}",
            doc.meta.get("school") and f"{doc.meta['school']}",
            doc.meta.get("major") and f"{doc.meta['major']}",
            doc.meta.get("class_name") and f"{doc.meta['class_name']}",
        ) if x)
        cover.placeholders[1].text = subtitle

        for sec in _sections(doc):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec["heading"]
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = []
            for text in sec["paragraphs"]:
                bullets.extend(_bullets(text, max_len=72))
                if len(bullets) >= 5:
                    break
            for i, text in enumerate(bullets[:5]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = text
                for run in para.runs:
                    run.font.size = Pt(18)

        if doc.annotations:
            chunks = [doc.annotations[i:i + 6] for i in range(0, len(doc.annotations), 6)]
            for page, lines in enumerate(chunks, 1):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "批注说明" + (f"（{page}/{len(chunks)}）" if len(chunks) > 1 else "")
                body = slide.placeholders[1].text_frame
                body.clear()
                for i, line in enumerate(lines):
                    para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    para.text = line
                    for run in para.runs:
                        run.font.size = Pt(12)

        refs = doc.references
        if refs:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "参考文献"
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, ref in enumerate(refs[:8]):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = ref[:80]
                for run in para.runs:
                    run.font.size = Pt(14)
        prs.save(path)

# 版本: v2.2.0 (2026-08-18) 更新: 批注说明导出
